import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def create_pptx_from_md(md_path: Path, output_path: Path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide divider
    slides_md = re.split(r'\n---\n+', content)
    
    prs = Presentation()
    
    # We will use simple slide layouts
    # 0 = Title Slide, 1 = Title and Content, 5 = Title Only, 6 = Blank
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    for i, slide_text in enumerate(slides_md):
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        lines = slide_text.split('\n')
        title = ""
        bullets = []
        texts = []
        image_name = None
        
        for line in lines:
            line = line.strip()
            if line.startswith('## '):
                title = line.replace('## ', '')
            elif line.startswith('# '):
                title = line.replace('# ', '')
            elif line.startswith('- '):
                bullets.append(line.replace('- ', ''))
            elif line.endswith('.png`') or line.endswith('.png'):
                m = re.search(r'([a-zA-Z0-9_]+\.png)', line)
                if m:
                    image_name = m.group(1)
            elif line != "" and not line.startswith('Visual:') and not line.startswith('Subtitle:'):
                texts.append(line)
                
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            title_shape.text = title if title else "Presentation"
            subtitle_shape.text = "\n".join(texts[:3])
            continue
            
        # Regular slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        if title:
            # clean up "Slide X - " prefix
            title = re.sub(r'Slide \d+ [—-] ', '', title)
            title_shape.text = title
            
        tf = body_shape.text_frame
        for b_idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if b_idx > 0 or texts else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            
        for t_idx, txt in enumerate(texts):
            if t_idx == 0 and not bullets:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.level = 0
            
        # Add image if referenced
        if image_name:
            img_path = Path("final_documentation/edge_opt_figures") / image_name
            if img_path.exists():
                # Add image to right side
                left = Inches(5.0)
                top = Inches(1.5)
                height = Inches(4.5)
                slide.shapes.add_picture(str(img_path), left, top, height=height)

    prs.save(str(output_path))
    print(f"Created {output_path}")

if __name__ == "__main__":
    create_pptx_from_md(
        Path("final_documentation/scientific_deck_v1.md"),
        Path("final_documentation/scientific_deck_v1.pptx")
    )
    create_pptx_from_md(
        Path("final_documentation/pitch_deck_v1.md"),
        Path("final_documentation/pitch_deck_v1.pptx")
    )
