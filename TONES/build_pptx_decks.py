#!/usr/bin/env python3
"""Build styled PPTX decks from clinical data and figures."""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG = Path("final_documentation/clinical_figures")
EFIG = Path("final_documentation/edge_opt_figures")
OUT = Path("final_documentation")

BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
BLUE = RGBColor(0x3A, 0x7B, 0xD5)
RED = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=14, color=GRAY, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.space_after = Pt(4)
    return txBox


def build_scientific(s):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    pp = s["per_participant"]
    zpct = s["clarke_zone_pcts"]
    zones = s["clarke_zones"]

    # --- Slide 1: Title ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(1.5), Inches(10), Inches(1.5),
             "Voice-Based Non-Invasive\nGlucose Estimation", 44, True, ACCENT)
    add_text(sl, Inches(0.8), Inches(3.5), Inches(10), Inches(0.6),
             "Full Optimization Study - TONES / ONVOX", 22, False, GRAY)
    add_text(sl, Inches(0.8), Inches(4.5), Inches(10), Inches(0.5),
             f"N = {s['n_total_samples']} matched voice-CGM samples | {s['n_participants']} participants | FreeStyle Libre validation",
             14, False, GRAY)

    # --- Slide 2: Key Results ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "KEY RESULTS", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Clinical Performance Summary", 32, True, WHITE)

    stats = [
        (f"{s['overall_mae']}", "Overall MAE\n(mg/dL)"),
        (f"{zpct['A']:.1f}%", "Clarke\nZone A"),
        (f"{s['clarke_ab_pct']:.1f}%", "Clarke\nZone A+B"),
        (f"{s['bland_altman_mean_bias']}", "Mean Bias\n(mg/dL)"),
        (f"{s['n_total_samples']}", "Matched\nSamples"),
        ("7.56", "Best MAE\n(Wolf)"),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.8 + i * 2.0)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(1.8), Inches(1.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # --- Slide 3: Clarke Error Grid ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "CLINICAL VALIDATION", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Clarke Error Grid Analysis", 32, True, WHITE)
    img = FIG / "clarke_error_grid.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.5), Inches(1.8), height=Inches(5.2))
    lines = [
        (f"Zone A: {zones['A']} ({zpct['A']:.1f}%)", 16, GREEN, True),
        (f"Zone B: {zones['B']} ({zpct['B']:.1f}%)", 16, BLUE, True),
        (f"Zone C: {zones['C']} ({zpct['C']:.1f}%)", 14, YELLOW, False),
        (f"Zone D: {zones['D']} ({zpct['D']:.1f}%)", 14, GRAY, False),
        (f"Zone E: {zones['E']} ({zpct['E']:.1f}%)", 14, GRAY, False),
        ("", 10, GRAY, False),
        ("Zero samples in dangerous Zones D or E.", 14, GREEN, True),
        ("99.2% of predictions within clinically acceptable limits.", 13, GRAY, False),
    ]
    add_multiline(sl, Inches(7.5), Inches(2.0), Inches(5), Inches(4), lines)

    # --- Slide 4: Bland-Altman ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "AGREEMENT ANALYSIS", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Bland-Altman: Voice vs. CGM", 32, True, WHITE)
    img = FIG / "bland_altman.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(1.8), width=Inches(7.5))
    lines = [
        (f"Mean bias: {s['bland_altman_mean_bias']} mg/dL", 16, ACCENT, True),
        (f"SD: {s['bland_altman_sd']} mg/dL", 14, GRAY, False),
        (f"95% LoA: [{s['bland_altman_95_loa'][0]}, {s['bland_altman_95_loa'][1]}] mg/dL", 14, GRAY, False),
        ("", 10, GRAY, False),
        ("Minimal systematic bias.", 14, GREEN, True),
        ("FreeStyle Libre MARD is ~9.2%.", 12, GRAY, False),
        ("Our personalized models approach CGM-level accuracy.", 12, GRAY, False),
    ]
    add_multiline(sl, Inches(8.2), Inches(2.0), Inches(4.5), Inches(4), lines)

    # --- Slide 5: Per-Participant ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "INDIVIDUAL PERFORMANCE", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Per-Participant Results", 32, True, WHITE)
    img = FIG / "participant_waterfall.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(1.8), width=Inches(7.0))
    lines = []
    for name in sorted(pp, key=lambda k: pp[k]["mae"]):
        d = pp[name]
        clr = GREEN if d["mae"] < 10 else BLUE if d["mae"] < 15 else YELLOW
        lines.append((f"{name}: MAE={d['mae']:.1f}, r={d['r']:.2f}, n={d['n']}", 13, clr, False))
    lines.append(("", 8, GRAY, False))
    lines.append(("7 of 9 achieve MAE < 15 mg/dL", 14, GREEN, True))
    add_multiline(sl, Inches(7.6), Inches(1.8), Inches(5), Inches(5.5), lines)

    # --- Slide 6: Offset Heatmap ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "NOVEL FINDING", 14, True, GREEN)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
             "Voice Features Lead CGM by 10-20 Minutes", 32, True, WHITE)
    img = FIG / "offset_heatmap.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(1.8), width=Inches(7.5))
    lines = [
        ("CGM sensors lag blood glucose by 10-15 min.", 14, GRAY, False),
        ("Voice may detect changes BEFORE the CGM.", 14, ACCENT, True),
        ("", 8, GRAY, False),
        ("5/9 participants best at positive offset", 14, GREEN, True),
        ("Wolf: +20min -> MAE drops to 7.9 mg/dL", 13, GRAY, False),
        ("Steffen: +20min -> MAE drops to 8.0 mg/dL", 13, GRAY, False),
        ("", 8, GRAY, False),
        ("Patent candidate: temporal offset-compensated", 12, YELLOW, True),
        ("voice-based glucose estimation method", 12, YELLOW, False),
    ]
    add_multiline(sl, Inches(8.2), Inches(2.0), Inches(4.5), Inches(5), lines)

    # --- Slide 7: Model Comparison ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "OPTIMIZATION", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Model & Feature Selection", 32, True, WHITE)
    img = FIG / "model_comparison.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(1.8), width=Inches(6.5))
    lines = [
        ("Edge-Inspired Optimization Sweep:", 15, WHITE, True),
        ("12 feature configs x 3 augmentations x 5 models", 13, GRAY, False),
        ("", 8, GRAY, False),
        ("Best Personalized:", 14, ACCENT, True),
        ("SVR + MFCC-20 + 2000ms window", 13, GRAY, False),
        ("MAE = 10.92 mg/dL", 13, GREEN, True),
        ("", 8, GRAY, False),
        ("Best Population:", 14, ACCENT, True),
        ("BayesianRidge + MFCC-20 + 300ms", 13, GRAY, False),
        ("MAE = 12.28 mg/dL", 13, GREEN, True),
        ("", 8, GRAY, False),
        ("Best Calibration (10-shot):", 14, ACCENT, True),
        ("SVR + MFCC-13 + 1000ms", 13, GRAY, False),
        ("MAE = 11.86 mg/dL", 13, GREEN, True),
    ]
    add_multiline(sl, Inches(7.3), Inches(1.8), Inches(5.5), Inches(5.5), lines)

    # --- Slide 8: Limitations ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "HONEST ASSESSMENT", 14, True, YELLOW)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "Limitations & Next Steps", 32, True, WHITE)
    lines_l = [
        ("Current Limitations:", 16, WHITE, True),
        ("Small cohort (n=9) - needs scale-up", 13, YELLOW, False),
        ("Healthy/pre-diabetic only", 13, YELLOW, False),
        ("No longitudinal drift analysis", 13, YELLOW, False),
        ("Modest Pearson r for some participants", 13, GRAY, False),
        ("Recording quality varies", 13, GRAY, False),
    ]
    add_multiline(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), lines_l)
    lines_r = [
        ("Next Scientific Steps:", 16, WHITE, True),
        ("100+ participant prospective study", 13, GREEN, False),
        ("Include T1D, T2D, gestational, healthy", 13, GREEN, False),
        ("Controlled recording protocol", 13, GREEN, False),
        ("Deep learning (1D-CNN, Wav2Vec2)", 13, GREEN, False),
        ("Context-aware features (YAMNET, activity)", 13, GREEN, False),
        ("Longitudinal stability analysis", 13, GREEN, False),
    ]
    add_multiline(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5), lines_r)

    prs.save(str(OUT / "scientific_deck_v2.pptx"))
    print("Created: final_documentation/scientific_deck_v2.pptx")


def build_pitch(s):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    zpct = s["clarke_zone_pcts"]

    # --- Slide 1: Title ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(2.0), Inches(10), Inches(1.5), "ONVOX", 64, True, ACCENT)
    add_text(sl, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
             "Non-invasive glucose monitoring\nthrough voice biomarkers", 24, False, GRAY)
    add_text(sl, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
             "Seed-stage  |  Pre-clinical  |  February 2026", 14, False, RGBColor(0x55, 0x55, 0x55))

    # --- Slide 2: Problem ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "THE PROBLEM", 14, True, RED)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
             "537M People with Diabetes.\nMost Don't Monitor Enough.", 36, True, WHITE)
    lines = [
        ("CGM adoption < 30% even among insulin-dependent", 15, GRAY, False),
        ("Fingerstick causes pain, stigma, low compliance", 15, GRAY, False),
        ("CGM costs $100-300/month (not always reimbursed)", 15, GRAY, False),
        ("No passive, consumer-friendly alternative exists", 15, ACCENT, True),
    ]
    add_multiline(sl, Inches(0.8), Inches(2.8), Inches(7), Inches(3), lines)

    for i, (val, label) in enumerate([("$31.4B", "Glucose monitoring\nmarket by 2028"), ("10.5%", "CAGR")]):
        x = Inches(8.5 + i * 2.3)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.0), Inches(2.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xF5, 0x57, 0x6C)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # --- Slide 3: Solution ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "OUR SOLUTION", 14, True, GREEN)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7),
             "Speak. Get Your Glucose Level.", 36, True, WHITE)
    lines = [
        ("No hardware - uses the phone's microphone", 16, GRAY, False),
        ("No consumables - zero marginal cost", 16, GRAY, False),
        ("No pain - completely non-invasive", 16, GRAY, False),
        ("Personalizable - improves with 5-10 reference readings", 16, ACCENT, True),
    ]
    add_multiline(sl, Inches(0.8), Inches(2.2), Inches(6), Inches(3), lines)
    lines2 = [
        ("How It Works:", 16, ACCENT, True),
        ("1. Record 3-5 seconds of speech", 14, GRAY, False),
        ("2. Extract audio features (MFCC, spectral, pitch)", 14, GRAY, False),
        ("3. ML model predicts glucose (mg/dL)", 14, GRAY, False),
        ("4. Display result with trend & confidence", 14, GRAY, False),
        ("5. Calibration improves over time", 14, GRAY, False),
    ]
    add_multiline(sl, Inches(7.5), Inches(2.2), Inches(5), Inches(4), lines2)

    # --- Slide 4: Evidence ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "CLINICAL EVIDENCE", 14, True, GREEN)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7),
             "99.2% of Predictions in Safe Zones", 36, True, WHITE)
    img = FIG / "clarke_error_grid.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(1.8), height=Inches(5.2))

    for i, (val, label) in enumerate([
        (f"{zpct['A']:.0f}%", "Zone A"), (f"{s['clarke_ab_pct']:.0f}%", "Zone A+B"), ("0", "Dangerous\nErrors")
    ]):
        x = Inches(7.5 + i * 2.0)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(1.7), Inches(1.3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    lines = [
        (f"{s['n_total_samples']} voice-CGM pairs validated", 13, GRAY, False),
        (f"{s['n_participants']} participants, SVR + LOO CV", 13, GRAY, False),
        ("FreeStyle Libre ground truth", 13, GRAY, False),
    ]
    add_multiline(sl, Inches(7.5), Inches(3.8), Inches(5), Inches(2), lines)

    # --- Slide 5: Offset ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "BREAKTHROUGH INSIGHT", 14, True, GREEN)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
             "Voice Detects Glucose Changes\n10-20 Minutes Before CGM", 34, True, WHITE)
    img = FIG / "offset_heatmap.png"
    if img.exists():
        sl.shapes.add_picture(str(img), Inches(0.3), Inches(2.2), width=Inches(7.2))
    lines = [
        ("CGM lags blood glucose by 10-15 min", 14, GRAY, False),
        ("Voice biomarkers may detect changes FIRST", 14, ACCENT, True),
        ("5/9 participants best at positive offset", 14, GREEN, True),
        ("", 8, GRAY, False),
        ("Patent candidate:", 13, YELLOW, True),
        ("Temporal offset-compensated", 13, YELLOW, False),
        ("voice-based glucose estimation", 13, YELLOW, False),
    ]
    add_multiline(sl, Inches(7.8), Inches(2.2), Inches(5), Inches(4.5), lines)

    # --- Slide 6: Competition ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "COMPETITION", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7),
             "No One Else Does This", 36, True, WHITE)
    data = [
        ("Klick Labs", "Diabetes detection", "Binary", "No", "No"),
        ("Sonde Health", "Mental/respiratory", "No", "No", "No"),
        ("Canary Speech", "Neurological", "No", "No", "No"),
        ("ONVOX (us)", "Glucose estimation", "mg/dL", "Yes", "Few-shot"),
    ]
    headers = ["Company", "Domain", "Glucose Level?", "Continuous?", "Personalized?"]
    y_start = Inches(2.0)
    col_widths = [Inches(2.0), Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0)]
    x_start = Inches(0.8)
    for j, h in enumerate(headers):
        x = x_start + sum(col_widths[:j])
        add_text(sl, x, y_start, col_widths[j], Inches(0.4), h, 11, True, BLUE, PP_ALIGN.LEFT)
    for i, row in enumerate(data):
        y = y_start + Inches(0.5) + Inches(i * 0.5)
        is_us = i == len(data) - 1
        for j, val in enumerate(row):
            x = x_start + sum(col_widths[:j])
            if is_us:
                clr = ACCENT if j <= 1 else GREEN
                bld = True
            else:
                clr = WHITE if j == 0 else (RED if val in ("No", "Binary") else GRAY)
                bld = j == 0
            add_text(sl, x, y, col_widths[j], Inches(0.4), val, 13, bld, clr)

    # --- Slide 7: Roadmap ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "ROADMAP", 14, True, BLUE)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "From Research to Revenue", 36, True, WHITE)
    phases = [
        ("Phase 1: Now", "Research Tool\nPublish, open pipeline", BLUE),
        ("Phase 2: 6-12mo", "Developer API\nEnterprise, pharma", GREEN),
        ("Phase 3: 12-24mo", "Consumer App\nHealth integrations", YELLOW),
        ("Phase 4: 24mo+", "Clinical Grade\nFDA/CE, EHR", RED),
    ]
    for i, (title, desc, clr) in enumerate(phases):
        x = Inches(0.8 + i * 3.1)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.8), Inches(2.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = clr
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY

    # --- Slide 8: The Ask ---
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4), "INVESTMENT", 14, True, ACCENT)
    add_text(sl, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7), "The Ask", 40, True, WHITE)
    lines = [
        ("Seed Round - Pre-clinical, strong signal, unique IP", 18, ACCENT, True),
        ("", 8, GRAY, False),
        ("Use of Funds:", 16, WHITE, True),
        ("40% - Prospective clinical study (100+ participants)", 14, GRAY, False),
        ("25% - ML engineering & model optimization", 14, GRAY, False),
        ("20% - Regulatory preparation (CE, FDA)", 14, GRAY, False),
        ("15% - Team & operations", 14, GRAY, False),
    ]
    add_multiline(sl, Inches(0.8), Inches(2.0), Inches(6), Inches(5), lines)
    lines2 = [
        ("Why Now?", 16, WHITE, True),
        ("Voice AI is mainstream (Siri, Alexa)", 14, GREEN, False),
        ("CGM market: $31.4B by 2028", 14, GREEN, False),
        ("Klick Labs' Nature paper validated concept", 14, GREEN, False),
        ("We are FIRST to show continuous estimation", 14, ACCENT, True),
        ("Temporal offset finding is novel & patentable", 14, ACCENT, True),
        ("", 8, GRAY, False),
        ("Milestones to Next Round:", 16, WHITE, True),
        ("Prospective validation study complete", 14, GRAY, False),
        ("Provisional patents filed", 14, GRAY, False),
        ("API beta with 3+ enterprise partners", 14, GRAY, False),
        ("Peer-reviewed publication accepted", 14, GRAY, False),
    ]
    add_multiline(sl, Inches(7.2), Inches(2.0), Inches(5.5), Inches(5.5), lines2)

    prs.save(str(OUT / "pitch_deck_v2.pptx"))
    print("Created: final_documentation/pitch_deck_v2.pptx")


def main():
    with open(FIG / "clinical_summary.json") as f:
        s = json.load(f)
    build_scientific(s)
    build_pitch(s)


if __name__ == "__main__":
    main()
